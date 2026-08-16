"""LibreOffice가 없어 렌더링 QA를 못 하므로 좌표로 검사한다.

  · 슬라이드 밖으로 나가는 도형/이미지
  · 좌우 0.62" 여백 침범
  · 결론 줄(y=6.62) 침범
  · 서로 겹치는 블록 (포함 관계는 배경 카드이므로 제외)
"""
from pptx import Presentation
from pptx.util import Emu

W, H, M, CONCL = 13.333, 7.5, 0.62, 6.62
p = Presentation("/Users/gimsiu/Downloads/suni c 14조/docs/발표_추세분석_판정.pptx")
E = lambda v: Emu(v).inches
bad = 0

for i, s in enumerate(p.slides, 1):
    items = []
    for sh in s.shapes:
        if sh.left is None:
            continue
        x, y = E(sh.left), E(sh.top)
        w, h = E(sh.width or 0), E(sh.height or 0)
        txt = (sh.text_frame.text[:26].replace("\n", "⏎") if sh.has_text_frame else
               f"<{sh.shape_type}>")
        items.append((x, y, w, h, txt))
        if x < -0.01 or y < -0.01 or x + w > W + 0.01 or y + h > H + 0.01:
            print(f"[{i}] 슬라이드 밖  {txt!r}  ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
            bad += 1
        if x < M - 0.01 or x + w > W - M + 0.01:
            print(f"[{i}] 여백 침범   {txt!r}  x={x:.2f}..{x + w:.2f}")
            bad += 1
        # 결론 줄 위 콘텐츠가 CONCL을 넘어 내려오면 겹친다 (결론/각주 자체는 제외)
        if y < CONCL - 0.02 and y + h > CONCL + 0.02:
            print(f"[{i}] 결론줄 침범 {txt!r}  y={y:.2f}..{y + h:.2f}")
            bad += 1

    for a in range(len(items)):
        for b in range(a + 1, len(items)):
            ax, ay, aw, ah, at = items[a]
            bx, by, bw, bh, bt = items[b]
            ox = min(ax + aw, bx + bw) - max(ax, bx)
            oy = min(ay + ah, by + bh) - max(ay, by)
            if ox <= 0.02 or oy <= 0.02:
                continue
            # 한쪽이 다른 쪽을 거의 담고 있으면 배경 카드 — 정상
            if ox >= min(aw, bw) - 0.03 and oy >= min(ah, bh) - 0.03:
                continue
            print(f"[{i}] 겹침 {ox:.2f}x{oy:.2f}  {at!r}  ×  {bt!r}")
            bad += 1

print(f"\n슬라이드 {len(p.slides)}장 · 문제 {bad}건")
